"""Build the six-slide SIH idea deck without changing the supplied theme/master."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "SIH2026-IDEA-Presentation-Format.pptx"
OUTPUT = ROOT / "ReliefChain_SIH2026_Idea_Presentation.pptx"

# Replace these three portal-controlled values before final submission.
PROBLEM_STATEMENT_ID = "[ENTER PS ID]"
OFFICIAL_THEME = "[ENTER OFFICIAL SIH THEME]"
TEAM_ID = "[ENTER TEAM ID]"
REGISTERED_TEAM_NAME = "RELIEFCHAIN [VERIFY PORTAL NAME]"

ORANGE = RGBColor(242, 112, 18)
GREEN = RGBColor(19, 138, 67)
SLATE = RGBColor(78, 98, 108)
DARK = RGBColor(35, 49, 56)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(244, 246, 247)
PALE_GREEN = RGBColor(234, 246, 238)
PALE_ORANGE = RGBColor(255, 243, 230)
PALE_BLUE = RGBColor(235, 242, 246)
RED = RGBColor(177, 43, 43)


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    relationship_id = slide_id.rId
    prs.part.drop_rel(relationship_id)
    del prs.slides._sldIdLst[index]


def set_text(shape, text: str, *, size: float = 12, bold: bool = False, color: RGBColor = DARK,
             align: PP_ALIGN = PP_ALIGN.CENTER, font_name: str = "Arial") -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text: str, *, size: float = 15, bold: bool = False,
             color: RGBColor = DARK, align: PP_ALIGN = PP_ALIGN.LEFT,
             valign: MSO_ANCHOR = MSO_ANCHOR.TOP, font_name: str = "Arial"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_multiline(slide, x, y, w, h, lines: Iterable[tuple[str, bool]], *, size: float = 13,
                  color: RGBColor = DARK, bullet: bool = False, spacing: float = 1.05):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.09)
    frame.margin_right = Inches(0.09)
    frame.margin_top = Inches(0.06)
    frame.margin_bottom = Inches(0.04)
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, (text, bold) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(size * 0.42 * spacing)
        run = paragraph.add_run()
        run.text = ("• " if bullet else "") + text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def add_card(slide, x, y, w, h, title: str, lines: list[str], *, accent: RGBColor = GREEN,
             fill: RGBColor = LIGHT, title_size: float = 15, body_size: float = 11.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = accent
    card.line.width = Pt(1.4)
    add_text(slide, x + 0.15, y + 0.12, w - 0.3, 0.34, title, size=title_size, bold=True, color=accent)
    add_multiline(slide, x + 0.12, y + 0.55, w - 0.24, h - 0.65,
                  [(line, False) for line in lines], size=body_size, bullet=True)
    return card


def add_pill(slide, x, y, w, text: str, color: RGBColor):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    set_text(pill, text, size=10.5, bold=True, color=WHITE)
    return pill


def add_flow_step(slide, x, y, w, title: str, subtitle: str, number: int, color: RGBColor):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.15))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.08), Inches(y + 0.13), Inches(0.38), Inches(0.38))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    set_text(badge, str(number), size=10, bold=True, color=WHITE)
    add_text(slide, x + 0.52, y + 0.10, w - 0.62, 0.37, title, size=12, bold=True, color=color)
    add_text(slide, x + 0.12, y + 0.53, w - 0.24, 0.5, subtitle, size=9.5, color=DARK, align=PP_ALIGN.CENTER)


def add_arrow(slide, x, y, w=0.34):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.28))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = SLATE
    arrow.line.fill.background()


def replace_team_badges(prs: Presentation) -> None:
    for slide in list(prs.slides)[1:]:
        for shape in slide.shapes:
            if shape.name.startswith("Oval"):
                set_text(shape, "RELIEFCHAIN", size=10, bold=True, color=DARK)


def clear_template_body(slide) -> None:
    for shape in list(slide.shapes):
        if shape.name == "TextBox 8":
            remove_shape(shape)


def add_notes(slide, text: str) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = text
    except AttributeError:
        pass


def build() -> None:
    prs = Presentation(SOURCE)
    delete_slide(prs, 6)  # Important Instructions: rules explicitly permit deletion before submission.
    replace_team_badges(prs)

    # Slide 1 — Title page. Existing master artwork and SIH logo remain untouched.
    slide = prs.slides[0]
    subtitle = next(shape for shape in slide.shapes if shape.name == "Subtitle 3")
    subtitle.text_frame.clear()
    subtitle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = subtitle.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "RELIEFCHAIN"; r.font.name = "Garamond"; r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = GREEN
    p2 = subtitle.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = "Verify every rupee. Expose no beneficiary."; r2.font.name = "Arial"; r2.font.size = Pt(17); r2.font.bold = True; r2.font.color.rgb = SLATE

    details = next(shape for shape in slide.shapes if shape.name == "TextBox 9")
    frame = details.text_frame
    frame.clear(); frame.word_wrap = True; frame.margin_left = Inches(0.04); frame.margin_right = Inches(0.05)
    values = [
        ("Problem Statement ID", PROBLEM_STATEMENT_ID),
        ("Problem Statement Title", "Transparent, Privacy-Preserving Disaster Relief Fund Tracking"),
        ("Theme", OFFICIAL_THEME),
        ("PS Category", "Software"),
        ("Team ID", TEAM_ID),
        ("Team Name", REGISTERED_TEAM_NAME),
    ]
    for index, (label, value) in enumerate(values):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.space_after = Pt(7)
        label_run = para.add_run(); label_run.text = f"{label}: "; label_run.font.name = "Arial"; label_run.font.size = Pt(16); label_run.font.bold = True; label_run.font.color.rgb = SLATE
        value_run = para.add_run(); value_run.text = value; value_run.font.name = "Arial"; value_run.font.size = Pt(16); value_run.font.bold = False; value_run.font.color.rgb = DARK
    add_notes(slide, "Open with the trust problem: beneficiaries and auditors currently depend on the same intermediaries they are trying to verify. ReliefChain turns claims into attributable ledger evidence while keeping recipient identity private.")

    # Slide 2 — Idea title / solution / uniqueness.
    slide = prs.slides[1]; clear_template_body(slide)
    add_text(slide, 0.7, 1.05, 11.85, 0.24, "PROPOSED SOLUTION  •  HOW IT ADDRESSES THE PROBLEM  •  INNOVATION & UNIQUENESS", size=9.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.7, 1.32, 11.85, 0.43, "RELIEFCHAIN", size=25, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, 1.2, 1.72, 10.85, 0.45, "A permissioned evidence trail from fund source to family payout", size=16, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    flow = [
        ("SOURCE", "Govt • NGO • donor"),
        ("ALLOCATE", "District + scheme"),
        ("COMMIT", "HMAC beneficiary link"),
        ("DISBURSE", "Reserve → settle/fail"),
        ("VERIFY", "Public • family • audit"),
    ]
    x = 0.55
    for index, (title, subtitle_text) in enumerate(flow, 1):
        width = 2.13
        color = ORANGE if index in (1, 4) else GREEN
        add_flow_step(slide, x, 2.32, width, title, subtitle_text, index, color)
        if index < len(flow): add_arrow(slide, x + width + 0.05, 2.75)
        x += 2.48
    add_card(slide, 0.55, 3.78, 3.9, 2.35, "PROBLEM", [
        "Funds cross many opaque accounts",
        "Families cannot verify promised aid",
        "Audits are slow and paper-heavy",
    ], accent=ORANGE, fill=PALE_ORANGE)
    add_card(slide, 4.72, 3.78, 3.9, 2.35, "SOLUTION", [
        "Fabric records accepted fund transitions",
        "Role-specific API views preserve privacy",
        "Committed events drive reconciliation",
    ], accent=GREEN, fill=PALE_GREEN)
    add_card(slide, 8.89, 3.78, 3.9, 2.35, "UNIQUE VALUE", [
        "Trust shifts from a claim to shared evidence",
        "Public proof without public beneficiary identity",
        "Ledger–database mismatch becomes detectable",
    ], accent=SLATE, fill=PALE_BLUE)
    add_notes(slide, "Explain the novelty in one sentence: ReliefChain combines end-to-end fund-state evidence with audience-specific privacy, so transparency no longer means exposing beneficiaries.")

    # Slide 3 — Technical approach.
    slide = prs.slides[2]; clear_template_body(slide)
    technologies = [
        ("Next.js", ORANGE), ("Flutter", GREEN), ("NestJS", SLATE),
        ("PostgreSQL", SLATE), ("Hyperledger Fabric", GREEN), ("Docker", ORANGE),
    ]
    start_x = 0.6
    for label, color in technologies:
        width = 1.65 if label != "Hyperledger Fabric" else 2.25
        add_pill(slide, start_x, 1.27, width, label, color)
        start_x += width + 0.16
    add_text(slide, 0.7, 1.72, 11.8, 0.24, "TECHNOLOGIES USED  •  METHODOLOGY & IMPLEMENTATION FLOW", size=9.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    steps = [
        ("ACTORS", "Public dashboard\nOperator portal\nBeneficiary mobile\nAuditor API", ORANGE),
        ("API CONTROL", "JWT + RBAC\nOrg ownership\nEligibility + balances\nIdempotency", SLATE),
        ("PRIVACY", "HMAC identifier\nAES-GCM contacts\nHashed phone\nAllowlisted events", GREEN),
        ("LEDGER", "8 Fabric v1 writes\nStable error codes\nSigned MSP actor\nReadAsset + history", ORANGE),
        ("EVIDENCE", "Committed indexer\nPublic proof\nReconciliation\nCSV / JSON trail", GREEN),
    ]
    x = 0.55
    for index, (title, body, color) in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.02), Inches(2.18), Inches(2.78))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = color; box.line.width = Pt(1.5)
        add_text(slide, x + 0.1, 2.17, 1.98, 0.35, title, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_multiline(slide, x + 0.12, 2.65, 1.94, 1.9, [(line, False) for line in body.split("\n")], size=10.5, bullet=True)
        if index < len(steps) - 1: add_arrow(slide, x + 2.23, 3.18, 0.34)
        x += 2.53
    add_text(slide, 0.72, 5.08, 11.9, 0.35, "PRIVATE OPERATIONS", size=12, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.72, 5.45, 5.68, 0.7, "PostgreSQL\nEncrypted PII • sessions • jobs • attempts", size=11.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, 6.64, 5.45, 5.68, 0.7, "FABRIC CHANNEL\nTamper-evident state • timestamped safe events", size=11.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_notes(slide, "Walk left to right. The client never talks directly to a peer. The API enforces application authorization, Fabric enforces deterministic ledger rules, and only committed privacy-safe events feed audit evidence.")

    # Slide 4 — Feasibility and viability.
    slide = prs.slides[3]; clear_template_body(slide)
    add_text(slide, 0.7, 1.05, 11.8, 0.24, "FEASIBILITY ANALYSIS  •  POTENTIAL RISKS  •  MITIGATION STRATEGIES", size=9.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    add_card(slide, 0.6, 1.32, 3.7, 3.82, "FEASIBLE NOW", [
        "Working NestJS + PostgreSQL backend",
        "3-org Fabric network and deployed chaincode",
        "Restart-safe demo seed and payout worker",
        "Health, readiness, audit and proof APIs",
        "Contract, privacy and invariant test suites",
    ], accent=GREEN, fill=PALE_GREEN, body_size=11)
    add_card(slide, 4.55, 1.32, 3.7, 3.82, "KEY RISKS", [
        "Garbage-in: ledger cannot verify field truth",
        "Single-host peers are not institutionally independent",
        "Identity, OTP and payout providers are simulated",
        "Ledger commit and DB write are not one ACID action",
        "Key custody and certificate rotation need hardening",
    ], accent=ORANGE, fill=PALE_ORANGE, body_size=11)
    add_card(slide, 8.5, 1.32, 3.7, 3.82, "MITIGATION PATH", [
        "Regulated UIDAI/payment-provider attestations",
        "Independent peers + multi-org endorsement",
        "Outbox/recovery and projection reconciliation",
        "HSM-backed keys + Fabric CA lifecycle",
        "Phased district pilot with audit checkpoints",
    ], accent=SLATE, fill=PALE_BLUE, body_size=11)
    add_text(slide, 0.7, 5.35, 11.8, 0.35, "ROLLOUT", size=12, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    rollout = [("1  HACKATHON MVP", "Synthetic Assam data"), ("2  CONTROLLED PILOT", "One disaster • selected districts"), ("3  SCALE", "Independent institutions • real providers")]
    x = 0.72
    for index, (title, body) in enumerate(rollout):
        color = GREEN if index == 0 else ORANGE if index == 1 else SLATE
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.72), Inches(3.76), Inches(0.65))
        box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
        set_text(box, f"{title}\n{body}", size=10.5, bold=True, color=WHITE)
        x += 4.03
    add_notes(slide, "Be explicit that feasibility is demonstrated by a working stack, while viability requires independent institutions and real regulated providers. This honesty strengthens rather than weakens the proposal.")

    # Slide 5 — Impact and benefits.
    slide = prs.slides[4]; clear_template_body(slide)
    add_text(slide, 0.7, 1.05, 11.8, 0.24, "TARGET-AUDIENCE IMPACT  •  SOCIAL, ECONOMIC & OPERATIONAL BENEFITS", size=9.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    cards = [
        ("BENEFICIARIES", ["Private OTP access", "Own payout status", "Opaque proof reference"], GREEN, PALE_GREEN),
        ("PUBLIC", ["District/scheme totals", "No public recipient PII", "Proof lookup by reference"], ORANGE, PALE_ORANGE),
        ("AUDITORS", ["Searchable event trail", "Faster reconciliation", "Export + direct ledger check"], SLATE, PALE_BLUE),
        ("GOVT + NGOs", ["Attributable actions", "Balance/ownership rules", "Shared source of evidence"], GREEN, PALE_GREEN),
    ]
    x_positions = [0.55, 3.68, 6.81, 9.94]
    for x, (title, lines, accent, fill) in zip(x_positions, cards):
        add_card(slide, x, 1.35, 2.84, 3.18, title, lines, accent=accent, fill=fill, title_size=13, body_size=10.7)
    add_text(slide, 0.75, 4.85, 11.8, 0.35, "EXPECTED SYSTEM OUTCOMES", size=13, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    outcomes = [
        ("TRACEABLE", "Source → allocation → payout"),
        ("TAMPER-EVIDENT", "History cannot be silently rewritten"),
        ("PRIVACY-SAFE", "Evidence without identity exposure"),
        ("REUSABLE", "New disaster, district or scheme"),
    ]
    x = 0.68
    for index, (title, body) in enumerate(outcomes):
        color = ORANGE if index % 2 else GREEN
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(5.35), Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
        set_text(circle, str(index + 1), size=12, bold=True, color=WHITE)
        add_text(slide, x + 0.68, 5.28, 2.15, 0.27, title, size=11, bold=True, color=color)
        add_text(slide, x + 0.68, 5.58, 2.15, 0.52, body, size=9.5, color=DARK)
        x += 3.05
    add_notes(slide, "Frame impact by audience. Do not claim a percentage reduction without a field pilot; describe the observable outcomes the MVP already enables.")

    # Slide 6 — Research and references.
    slide = prs.slides[5]; clear_template_body(slide)
    add_text(slide, 0.7, 1.05, 11.8, 0.24, "RESEARCH WORK  •  OFFICIAL TECHNICAL & PRIVACY REFERENCES", size=9.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    references = [
        ("1", "Hyperledger Fabric — security model, identities and permissioned governance", "hyperledger-fabric.readthedocs.io/en/latest/security_model.html", GREEN),
        ("2", "Hyperledger Fabric — tamper-resistant ledger, world state and history", "hyperledger-fabric.readthedocs.io/en/release-2.5/ledger/ledger.html", ORANGE),
        ("3", "NIST — AES-GCM authenticated encryption and HMAC standards", "csrc.nist.gov/pubs/sp/800/38/d/final  •  csrc.nist.gov/pubs/fips/198-1/final", SLATE),
        ("4", "India privacy context — DPDP Act 2023; UIDAI Aadhaar Data Vault guidance", "meity.gov.in  •  uidai.gov.in/images/Circular-No.14_of-2025.pdf", GREEN),
    ]
    y = 1.32
    hyperlinks = [
        "https://hyperledger-fabric.readthedocs.io/en/latest/security_model.html",
        "https://hyperledger-fabric.readthedocs.io/en/release-2.5/ledger/ledger.html",
        "https://csrc.nist.gov/pubs/sp/800/38/d/final",
        "https://www.meity.gov.in/writereaddata/files/Digital%20Personal%20Data%20Protection%20Act%202023.pdf",
    ]
    for (number, title, link_text, color), hyperlink in zip(references, hyperlinks):
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.72), Inches(y + 0.08), Inches(0.48), Inches(0.48))
        badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background(); set_text(badge, number, size=11, bold=True, color=WHITE)
        add_text(slide, 1.38, y, 10.9, 0.32, title, size=12.5, bold=True, color=color)
        link_shape = add_text(slide, 1.38, y + 0.37, 10.9, 0.35, link_text, size=9.5, color=SLATE)
        link_shape.text_frame.paragraphs[0].runs[0].hyperlink.address = hyperlink
        y += 1.04
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(5.55), Inches(11.55), Inches(0.72))
    box.fill.solid(); box.fill.fore_color.rgb = PALE_GREEN; box.line.color.rgb = GREEN
    set_text(box, "MVP EVIDENCE  •  Fabric ReadAsset + GetHistory  •  committed-event index  •  contract/privacy/invariant tests", size=11.5, bold=True, color=GREEN)
    add_notes(slide, "References are official primary sources. Clarify that UIDAI integration is not claimed; the current MVP uses synthetic identifiers and cites UIDAI guidance as a production design constraint.")

    prs.core_properties.title = "ReliefChain — SIH 2026 Idea Presentation"
    prs.core_properties.subject = "Privacy-preserving, blockchain-verifiable disaster relief fund tracking"
    prs.core_properties.author = REGISTERED_TEAM_NAME
    prs.core_properties.comments = "Built from the supplied SIH 2026 template; slide master and theme retained."
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
